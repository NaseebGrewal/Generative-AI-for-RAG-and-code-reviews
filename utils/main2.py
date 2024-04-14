import os

import fitz

# import textract
from Azure_Cognitive_Search import (
    create_index,
    delete_index,
    generate_random_id,
    semantic_search,
    upload_to_index,
)
from docx import Document
from flask import Flask, redirect, render_template, request, session, url_for
from open_ai_response import response_from_open_ai
from pptx import Presentation
from Speech_to_text import get_text
from werkzeug.utils import secure_filename

app = Flask(__name__)
app.config["UPLOAD_FOLDER"] = "uploads/"
app.config["ALLOWED_EXTENSIONS"] = {"pdf", "txt", "pptx", "docx", "wav"}
app.config["ALLOWED_EXTENSIONS2"] = {"py"}

messages = [
    {
        "role": "system",
        "content": "You are a helpful AI assistant",
    }
]

messages2 = [
    {
        "role": "system",
        "content": """You are a python code review expert, who checks the code according to ruff, 
                an extremely fast python linter, by executing "ruff check ." command on the python code.
                """,
    }
]

chat_history = []
chat_history2 = []

app.secret_key = generate_random_id()
delete_index("test-index")


def allowed_file(filename):
    return (
        "." in filename
        and filename.rsplit(".", 1)[1].lower() in app.config["ALLOWED_EXTENSIONS"]
    )


def allowed_file2(filename):
    return (
        "." in filename
        and filename.rsplit(".", 1)[1].lower() in app.config["ALLOWED_EXTENSIONS2"]
    )


def extract_text(filepath):
    ext = filepath.rsplit(".", 1)[1].lower()
    chunks = []
    text = ""

    try:
        if ext == "pdf":
            with fitz.open(filepath) as doc:
                # text = []
                for page in doc:
                    chunks.append(page.get_text())
            # print(chunks)

        elif ext == "txt":
            with open(filepath, "r", encoding="utf-8") as file:
                text = file.read()
                chunks.append(text)

        elif ext == "wav":
            text = get_text(filepath)
            chunks.append(text)

        elif ext == "pptx":
            presentation = Presentation(filepath)
            for slide in presentation.slides:
                slide_text = ""
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text += shape.text + "\n"
                if slide_text:
                    chunks.append(slide_text.strip())

        elif ext == "docx":
            doc = Document(filepath)
            for para in doc.paragraphs:
                text = para.text
                if text:
                    chunks.append(text)

        # else:
        #     text = textract.process(filepath).decode("utf-8")
        #     chunks.append(text)

    except FileNotFoundError:
        return f"File not found: {filepath}", ext
    except IOError:
        return f"Could not read file: {filepath}", ext
    except Exception as e:
        return f"An error occurred while extracting text: {e}", ext

    return chunks, ext


@app.route("/chatbot", methods=["GET", "POST"])
def chatbot():
    global chat_history

    if request.method == "POST":
        file = request.files["file"]
        if file and allowed_file(file.filename):
            filename = secure_filename(file.filename)
            filepath = os.path.join(app.config["UPLOAD_FOLDER"], filename)
            file.save(filepath)
            # Save filename in session for later use
            session["uploaded_filename"] = filename
            chunks, extension = extract_text(filepath)
            index = "test-index"
            session["index"] = index
            if create_index(index):
                upload_to_index(index, chunks, filename, extension)

            session.modified = True
            return redirect(url_for("chatbot"))

        user_query = request.form["query"]
        ai_response = get_ai_response(user_query)
        # ai_response = ai_response.replace("\n", "<br>")
        # user_query = user_query.replace("\n", "<br>")
        chat_history.append({"user": user_query, "bot": ai_response})
        return redirect(
            url_for("chatbot")
        )  # This will redirect to a GET request after POST

    # messages = session.get('messages', [])
    return render_template("chatbot.html", chat_history=chat_history)


@app.route("/remove_file", methods=["POST"])
def remove_file():
    filename = session.get("uploaded_filename")
    if filename:
        filepath = os.path.join(app.config["UPLOAD_FOLDER"], filename)

        # Remove the file from the filesystem
        if os.path.exists(filepath):
            os.remove(filepath)
        index = session.get("index")
        delete_index(index)
        # Add logic to remove the file from Azure Cognitive Search index
        # This is a placeholder, you need to implement the actual logic
        # remove_from_index(index, filename)

        # Clear the filename from the session
        session.pop("uploaded_filename", None)
        session.pop("index", None)
        session.modified = True

    return redirect(url_for("chatbot"))


@app.route("/code_review", methods=["GET", "POST"])
def code_review():
    global chat_history2
    if request.method == "POST":
        # Check if a file is part of the form
        file = request.files.get("file")
        if file and allowed_file2(file.filename):
            # If a Python file is uploaded, read its content
            filename = secure_filename(file.filename)
            filepath = os.path.join(app.config["UPLOAD_FOLDER"], filename)
            file.save(filepath)
            with open(filepath, "r", encoding="utf-8") as f:
                user_query = f.read()
            # Remove the file after reading its content if you wish
            os.remove(filepath)
        else:
            # If no file is uploaded, take the code from the textarea
            user_query = request.form["query"]

        ai_response = get_ai_response_code_review(user_query)
        # ai_response = ai_response.replace("\n", "<br>")
        # user_query = user_query.replace("\n", "<br>")
        chat_history2.append({"user": user_query, "bot": ai_response})
        session.modified = True
        return redirect(
            url_for("code_review")
        )  # This will redirect to a GET request after POST

    # messages = session.get('messages', [])

    return render_template("code_review.html", chat_history=chat_history2)


@app.route("/reset_chatbot", methods=["POST"])
def reset_chatbot():
    global chat_history
    chat_history = []
    # Reset the chat history
    # Ensure the session is marked as modified so it gets saved
    return redirect(url_for("chatbot"))  # Redirect back to the code review page


@app.route("/reset_code_review", methods=["POST"])
def reset_code_review():
    global chat_history2
    chat_history2 = []  # Reset the chat history
    # Ensure the session is marked as modified so it gets saved
    return redirect(url_for("code_review"))  # Redirect back to the code review page


def get_ai_response(query):
    """
    This function returns a response from a generative AI model.
    """
    # Retrieve messages value from global scope
    global messages

    if "index" in session:
        data = semantic_search(query, "test-index")

        if data:
            global messages
            context = [dct["content"] for dct in data]
            context = " ".join(context)
            # Construct the prompt with the user's query
            prompt = f"""Your task is to provide an answer to the user_query based on the given file content.
                        user_query: ```{query}```
                        Given file content: ```{context}```
                    """

            # Append the user's query to the messages
            messages.append({"role": "user", "content": prompt})

            text_to_add = "Most relevant "
            ext = data[0]["extension"]
            if ext == "pdf":
                text_to_add += "Page"
            elif ext == "docx":
                text_to_add += "Paragraph"
            elif ext == "pptx":
                text_to_add += "Slide"
            else:
                text_to_add += "Chunk"

            page_para_slide = data[0]["page_para_slide"]
            text_to_add += f" : {page_para_slide}\nFile Type: {ext}\n"

            # Get the response from the AI model
            response = response_from_open_ai(messages)
            response = text_to_add + response
        else:
            response = f"""I am sorry, but there is nothing regarding "{query}" is mentioned in the document. Please ask something else."""

    else:
        messages.append({"role": "user", "content": query})
        response = response_from_open_ai(messages)
        text_to_add = "Since you haven't uploaded a file yet, I am only a general purpose chatbot, capable of answering on my trained knowledge.\n\n"
        response = text_to_add + response
    # Append the AI response to the messages
    messages.append({"role": "assistant", "content": response})

    return response


def get_ai_response_code_review(query):
    """
    This function returns a response from a generative AI model.
    """
    # Retrieve messages from session
    global messages2
    prompt = f"""Your task is to first find out whether the given user message is a valid python code or not.
            user message : ```{query}```
            if user message not a valid python code, then your task is to only respond "Please type valid python code".
            If it is a valid python code, then your task is to write two things respectively:
            1. Error messages when "ruff check ." command is executed on given python code.
            2. Correct Code, when all the error are fixed.
            """
    messages2.append({"role": "user", "content": prompt})

    # Get the response from the AI model
    response = response_from_open_ai(messages2)

    # Append the AI response to the messages
    messages2.append({"role": "assistant", "content": response})

    return response


if __name__ == "__main__":
    os.makedirs(app.config["UPLOAD_FOLDER"], exist_ok=True)
    app.run(debug=True)

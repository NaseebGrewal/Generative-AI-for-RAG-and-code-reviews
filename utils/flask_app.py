import os

import fitz

# import textract
from Azure_Cognitive_Search import (
    create_index,
    delete_index,
    generate_random_id,
    semantic_search,
    upload_to_index,
)
from code_review import examples
from docx import Document
from flask import Flask, escape, redirect, render_template, request, session, url_for
from open_ai_response import (
    response_from_open_ai,
    response_from_open_ai1,
    response_from_open_ai2,
)
from pptx import Presentation
from Speech_to_text import get_text
from werkzeug.utils import secure_filename

app = Flask(__name__)
app.config["UPLOAD_FOLDER"] = "uploads/"
app.config["ALLOWED_EXTENSIONS"] = {"pdf", "txt", "pptx", "docx", "wav"}
app.config["ALLOWED_EXTENSIONS2"] = {"py"}

messages = [
    {
        "role": "system",
        "content": "You are a helpful AI assistant",
    }
]

messages2 = [
    {
        "role": "system",
        "content": """"Welcome to the Python Code Ruff Error Analysis Service. \
Please submit a Python script, and the system will carry out the following tasks: 
Determine if the submitted text is a valid Python script.
If the text is a Python script, identify and list all 'ruff' errors present in the script.
Provide a corrected version of the script with all 'ruff' errors resolved.
Note that a repository of 'example Python scripts' with identified 'ruff' errors is available for the AI to utilize as a reference. \
If the input is not Python code, the system will respond with a specific message indicating the need for valid Python code.
""",
    }
]

chat_history = []
chat_history2 = []

app.secret_key = generate_random_id()
delete_index("test-index")


def allowed_file(filename):
    return (
        "." in filename
        and filename.rsplit(".", 1)[1].lower() in app.config["ALLOWED_EXTENSIONS"]
    )


def allowed_file2(filename):
    return (
        "." in filename
        and filename.rsplit(".", 1)[1].lower() in app.config["ALLOWED_EXTENSIONS2"]
    )


def extract_text(filepath):
    ext = filepath.rsplit(".", 1)[1].lower()
    chunks = []
    text = ""

    try:
        if ext == "pdf":
            with fitz.open(filepath) as doc:
                # text = []
                for page in doc:
                    chunks.append(page.get_text())
            # print(chunks)

        elif ext == "txt":
            with open(filepath, "r", encoding="utf-8") as file:
                text = file.read()
                chunks.append(text)

        elif ext == "wav":
            text = get_text(filepath)
            chunks.append(text)

        elif ext == "pptx":
            presentation = Presentation(filepath)
            for slide in presentation.slides:
                slide_text = ""
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text += shape.text + "\n"
                if slide_text:
                    chunks.append(slide_text.strip())

        elif ext == "docx":
            doc = Document(filepath)
            for para in doc.paragraphs:
                text = para.text
                if text:
                    chunks.append(text)

        # else:
        #     text = textract.process(filepath).decode("utf-8")
        #     chunks.append(text)

    except FileNotFoundError:
        return f"File not found: {filepath}", ext
    except IOError:
        return f"Could not read file: {filepath}", ext
    except Exception as e:
        return f"An error occurred while extracting text: {e}", ext

    return chunks, ext


@app.route("/chatbot", methods=["GET", "POST"])
def chatbot():
    global chat_history

    if request.method == "POST":
        file = request.files["file"]
        if file and allowed_file(file.filename):
            filename = secure_filename(file.filename)
            filepath = os.path.join(app.config["UPLOAD_FOLDER"], filename)
            file.save(filepath)
            # Save filename in session for later use
            session["uploaded_filename"] = filename
            chunks, extension = extract_text(filepath)
            index = "test-index"
            session["index"] = index
            if create_index(index):
                upload_to_index(index, chunks, filename, extension)

            session.modified = True
            return redirect(url_for("chatbot"))

        # user_query = request.form['query']
        user_query = escape(request.form["query"])  # Escape the user input
        ai_response = get_ai_response(user_query)
        # ai_response = ai_response.replace("\n", "<br>")
        # user_query = user_query.replace("\n", "<br>")
        chat_history.append({"user": user_query, "bot": ai_response})
        return redirect(
            url_for("chatbot")
        )  # This will redirect to a GET request after POST

    # messages = session.get('messages', [])
    return render_template("chatbot.html", chat_history=chat_history)


@app.route("/remove_file", methods=["POST"])
def remove_file():
    filename = session.get("uploaded_filename")
    if filename:
        filepath = os.path.join(app.config["UPLOAD_FOLDER"], filename)

        # Remove the file from the filesystem
        if os.path.exists(filepath):
            os.remove(filepath)
        index = session.get("index")
        delete_index(index)
        # Add logic to remove the file from Azure Cognitive Search index
        # This is a placeholder, you need to implement the actual logic
        # remove_from_index(index, filename)

        # Clear the filename from the session
        session.pop("uploaded_filename", None)
        session.pop("index", None)
        session.modified = True

    return redirect(url_for("chatbot"))


@app.route("/code_review", methods=["GET", "POST"])
def code_review():
    global chat_history2
    if request.method == "POST":
        # Check if a file is part of the form
        file = request.files.get("file")
        if file and allowed_file2(file.filename):
            # If a Python file is uploaded, read its content
            filename = secure_filename(file.filename)
            filepath = os.path.join(app.config["UPLOAD_FOLDER"], filename)
            file.save(filepath)
            with open(filepath, "r", encoding="utf-8") as f:
                user_query = escape(f.read())
            # Remove the file after reading its content if you wish
            os.remove(filepath)
        else:
            # If no file is uploaded, take the code from the textarea
            user_query = escape(request.form["query"])

        ai_response = get_ai_response_code_review(user_query)
        # ai_response = ai_response.replace("\n", "<br>")
        # user_query = user_query.replace("\n", "<br>")
        chat_history2.append({"user": user_query, "bot": ai_response})
        session.modified = True
        return redirect(
            url_for("code_review")
        )  # This will redirect to a GET request after POST

    # messages = session.get('messages', [])

    return render_template("code_review.html", chat_history=chat_history2)


@app.route("/reset_chatbot", methods=["POST"])
def reset_chatbot():
    global chat_history
    chat_history = []
    # Reset the chat history
    # Ensure the session is marked as modified so it gets saved
    return redirect(url_for("chatbot"))  # Redirect back to the code review page


@app.route("/reset_code_review", methods=["POST"])
def reset_code_review():
    global chat_history2
    chat_history2 = []  # Reset the chat history
    # Ensure the session is marked as modified so it gets saved
    return redirect(url_for("code_review"))  # Redirect back to the code review page


def get_ai_response(query):
    """
    This function returns a response from a generative AI model.
    """
    # Retrieve messages value from global scope
    global messages
    print(session)
    if "index" in session:
        # global messages
        messages = [
            {
                "role": "system",
                "content": "You are a helpful AI assistant",
            }
        ]

        data = semantic_search(query, "test-index")

        # print(data)

        if data:
            context = [dct["content"] for dct in data]
            context = " ".join(context)
            file_name = data[0]["filename"]
            context = f"File name: {file_name}\n" + context
            # Construct the prompt with the user's query
            prompt = f"""Instructions for the AI:
Your responses must be solely derived from the information contained within the provided text from the file content. \
If the information required to answer the user's query is not present in the file content, \
you must respond with the following statement: "Given query is not mentioned in the file."

User Query Template:
user_query: {query}

File Content Source Template:
Given file content: {context}
"""

            # Append the user's query to the messages
            messages.append({"role": "user", "content": prompt})

            text_to_add = "Most relevant "
            ext = data[0]["extension"]
            if ext == "pdf":
                text_to_add += "Page"
            elif ext == "docx":
                text_to_add += "Paragraph"
            elif ext == "pptx":
                text_to_add += "Slide"
            else:
                text_to_add += "Chunk"

            page_para_slide = data[0]["page_para_slide"]
            text_to_add += f" : {page_para_slide}\nFile Type: {ext}\n"

            response = response_from_open_ai(messages)
            messages.append({"role": "assistant", "content": response})
            # Get the response from the AI model
            if not response == "Given query is not mentioned in the file":
                response = text_to_add + response
            print(1)
        else:
            data = semantic_search("", "test-index")
            context = [dct["content"] for dct in data]
            context = " ".join(context)
            file_name = data[0]["filename"]
            context = f"File name: {file_name}\n" + context
            # Construct the prompt with the user's query
            prompt = f"""Instruction for AI:
When responding to a user's question, you must only use information contained within the provided file content. \
If the file content does not contain the information necessary to answer the user's question, \
you must reply with "Given query is not mentioned in the file."

User's Question Format:
user_query: {query}

Provided File Content Format:
Given file content: {context}
"""

            # Append the user's query to the messages
            messages.append({"role": "user", "content": prompt})
            response = response_from_open_ai(messages)
            messages.append({"role": "assistant", "content": response})

            print(2)
            # response = f"""I am sorry, but there is nothing regarding "{query}" is mentioned in the document. Please ask something else."""

    else:
        messages.append({"role": "user", "content": query})
        response = response_from_open_ai1(messages)
        messages.append({"role": "assistant", "content": response})
        text_to_add = "Since you haven't uploaded a file yet, I am a general purpose chatbot, capable of answering on my trained knowledge.\n\n"
        response = text_to_add + response
        print(3)
    # Append the AI response to the messages

    return response


def get_ai_response_code_review(query):
    """
    This function returns a response from a generative AI model.
    """
    # Retrieve messages from session
    global messages2
    prompt = f"""Below is the text provided by the user. You are to perform the following actions:
Check whether the text is Python code.
If the text is not Python code, reply with: 'I detected no python code. Please input python code.'
If the text is Python code:
Using the prepared 'example Python scripts' and their associated 'ruff' errors for reference, \
identify and list all 'ruff' errors found in the provided script.
Generate a corrected version of the script that addresses and fixes all identified 'ruff' errors.

Instructions for AI:
Strictly follow the instructions laid out above.
Reference the prepared 'example Python scripts' and their documented 'ruff' errors to aid in your analysis of the provided script.
Ensure that your identification of 'ruff' errors is clear, precise, and comprehensive.
Provide the corrected Python script, ensuring that it retains the original intent and functionality while being free of 'ruff' errors.

Example Python Scripts Repository:
{examples}

(These examples are for the AI's reference and should be used to inform the analysis of the provided script.)

User's Script Input:

{query} 
 
End of Input

Upon completion of the analysis, list the 'ruff' errors first, followed by the corrected Python script.
"""
    messages2.append({"role": "user", "content": prompt})

    # Get the response from the AI model
    response = response_from_open_ai2(messages2)

    # Append the AI response to the messages
    messages2.append({"role": "assistant", "content": response})

    return response


if __name__ == "__main__":
    os.makedirs(app.config["UPLOAD_FOLDER"], exist_ok=True)
    app.run(debug=True)

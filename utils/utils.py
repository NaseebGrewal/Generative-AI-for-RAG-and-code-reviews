import docx
import openai
import pptx
import PyPDF2
import speech_recognition as sr

# Set your OpenAI API Key
openai.api_key = "your_openai_api_key"


# GPT-3 function
def generate_gpt3_response(prompt):
    response = openai.Completion.create(
        engine="davinci-codex",
        prompt=prompt,
        max_tokens=150,
        n=1,
        stop=None,
        temperature=0.5,
    )
    return response.choices[0].text.strip()


# Extract text from PDF
def process_pdf(file):
    pdf_file = open(file, "rb")
    pdf_reader = PyPDF2.PdfFileReader(pdf_file)
    text = ""
    for page in range(pdf_reader.numPages):
        text += pdf_reader.getPage(page).extract_text()
    pdf_file.close()
    return text


# Extract text from Word document
def process_docx(file):
    doc = docx.Document(file)
    text = ""
    for paragraph in doc.paragraphs:
        text += paragraph.text + "\n"
    return text


# Extract text from PowerPoint
def process_pptx(file):
    ppt = pptx.Presentation(file)
    text = ""
    for slide in ppt.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text


# Extract text from audio
def process_audio(file):
    recognizer = sr.Recognizer()
    audio_file = sr.AudioFile(file)
    with audio_file as source:
        audio = recognizer.record(source)
    text = recognizer.recognize_google(audio)
    return text


# # Extract text from video
# def process_video(file):
#     video = VideoFileClip(file)
#     audio = video.audio
#     audio.write_audiofile("temp_audio.wav")
#     text = process_audio("temp_audio.wav")
#     return text


# Main function
def main():
    file = input("Enter the file path: ")
    file_type = file.split(".")[-1]

    if file_type == "pdf":
        text = process_pdf(file)
    elif file_type == "docx":
        text = process_docx(file)
    elif file_type == "pptx":
        text = process_pptx(file)
    elif file_type in ["wav", "mp3", "m4a"]:
        text = process_audio(file)
    # elif file_type in ["mp4", "avi", "mov"]:
    #     text = process_video(file)
    else:
        print("Unsupported file type.")
        return

    conversation = f"Document content:\n{text}\n\n"

    while True:
        user_input = input("You: ")
        if user_input.lower() == "quit":
            break
        conversation += f"You: {user_input}\n"
        prompt = f"{conversation}AI:"
        response = generate_gpt3_response(prompt)
        print(f"AI: {response}")
        conversation += f"AI: {response}\n"


if __name__ == "__main__":
    main()
